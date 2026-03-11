"""Generate the GTM deck (.pptx) from structured data sources.

Output: outputs/gtm_pharma_deck.pptx

Focused on: active/recruiting trials, top sponsors, what can be matched today,
what still requires chart review or site data.
"""

from datetime import date

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    raise SystemExit("python-pptx not installed. Run: uv sync")

from collections import Counter

from load_data import (
    load_trials, load_memberships, load_fields, load_criteria,
    load_not_evaluable, unique_trial_count, trials_per_condition,
    recruiting_trials, active_trials, condition_label, ensure_outputs,
    OUTPUTS,
)

# Colors
NAVY = RGBColor(0x0C, 0x1E, 0x3A)
TEAL = RGBColor(0x00, 0x89, 0x9B)
LIGHT_TEAL = RGBColor(0x00, 0xB4, 0xD8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFFWHITE = RGBColor(0xF8, 0xFA, 0xFC)
LIGHT_GRAY = RGBColor(0xE2, 0xE8, 0xF0)
MID_GRAY = RGBColor(0x64, 0x74, 0x8B)
DARK_GRAY = RGBColor(0x33, 0x40, 0x55)
GREEN = RGBColor(0x10, 0xB9, 0x81)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
RED_ACCENT = RGBColor(0xEF, 0x44, 0x44)
CARD_BG_DARK = RGBColor(0x14, 0x2B, 0x4D)


def add_shape(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, size=14,
             color=DARK_GRAY, bold=False, align=PP_ALIGN.LEFT, name="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = name
    p.alignment = align
    return tb


def add_stat(slide, x, y, number, label, color=TEAL):
    add_text(slide, Inches(x), Inches(y), Inches(2.2), Inches(0.7),
             number, size=42, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(x), Inches(y + 0.65), Inches(2.2), Inches(0.5),
             label, size=12, color=MID_GRAY, align=PP_ALIGN.CENTER)


def generate():
    trials = load_trials()
    memberships = load_memberships()
    fields = load_fields()
    criteria = load_criteria()
    ensure_outputs()

    n_unique = unique_trial_count(trials)
    n_recruiting = len(recruiting_trials(trials))
    n_active = len(active_trials(trials))
    cond_counts = trials_per_condition(memberships)

    # Top sponsors among recruiting trials
    rec_trials = recruiting_trials(trials)
    sponsor_counts = Counter(t["sponsor"] for t in rec_trials)

    # Direct pre-screening fields
    direct_fields = [f for f in fields if f["confidence_for_prescreening"] == "direct"]
    direct_criteria = [c for c in criteria if c["confidence"] == "direct"]
    not_eval_criteria = [c for c in criteria if c["confidence"] == "not_evaluable"]

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]

    # ── SLIDE 1: TITLE ──
    slide = prs.slides.add_slide(blank)
    add_shape(slide, Inches(0), Inches(0), prs.slide_width, prs.slide_height, NAVY)
    add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), TEAL)
    add_text(slide, Inches(1.2), Inches(1.8), Inches(10), Inches(1),
             "iDHEA", size=52, color=TEAL, bold=True)
    add_text(slide, Inches(1.2), Inches(2.8), Inches(10), Inches(1.2),
             "Clinical Trial Pre-Screening\nand Feasibility for Pharma",
             size=32, color=WHITE, bold=True)
    add_text(slide, Inches(1.2), Inches(4.3), Inches(10), Inches(0.5),
             "No new imaging capture needed for anatomy-only pre-screening "
             "across 64 optometry sites",
             size=16, color=LIGHT_TEAL)
    add_text(slide, Inches(1.2), Inches(5.5), Inches(4), Inches(0.4),
             f"{date.today().strftime('%B %Y')}  |  Confidential",
             size=14, color=MID_GRAY)

    # ── SLIDE 2: THE PROBLEM ──
    slide = prs.slides.add_slide(blank)
    add_shape(slide, Inches(0), Inches(0), prs.slide_width, prs.slide_height, OFFWHITE)
    add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), TEAL)
    add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.5),
             "THE PROBLEM", size=14, color=TEAL, bold=True)
    add_text(slide, Inches(0.8), Inches(0.85), Inches(11), Inches(0.7),
             "Clinical trial enrollment is the #1 bottleneck in drug development",
             size=28, color=NAVY, bold=True)
    add_stat(slide, 0.8, 2.0, "$6,500\u2013$12K", "Cost per enrolled\nophthalmology patient")
    add_stat(slide, 3.8, 2.0, "85%+", "Screen failure rate\nin retinal trials")
    add_stat(slide, 6.8, 2.0, "6\u201312 mo", "Typical enrollment\ndelay per trial",
             color=RED_ACCENT)

    # Problem cards
    for i, (title, body, color) in enumerate([
        ("For Pharma Sponsors",
         "Every month of delay costs $600K\u2013$8M for a blockbuster drug.",
         RED_ACCENT),
        ("For CROs",
         "Site selection relies on surveys, not patient-level data.",
         AMBER),
        ("For Trial Sites",
         "Optometrists see eligible patients daily but cannot connect them to trials.",
         TEAL),
    ]):
        x = 0.8 + i * 4.0
        card = add_shape(slide, Inches(x), Inches(3.8), Inches(3.6), Inches(1.6), WHITE)
        add_shape(slide, Inches(x), Inches(3.8), Inches(0.06), Inches(1.6), color)
        add_text(slide, Inches(x + 0.25), Inches(3.9), Inches(3.1), Inches(0.35),
                 title, size=14, color=NAVY, bold=True)
        add_text(slide, Inches(x + 0.25), Inches(4.3), Inches(3.1), Inches(0.9),
                 body, size=12, color=MID_GRAY)

    # ── SLIDE 3: THE iDHEA ASSET ──
    slide = prs.slides.add_slide(blank)
    add_shape(slide, Inches(0), Inches(0), prs.slide_width, prs.slide_height, NAVY)
    add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), TEAL)
    add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.5),
             "THE iDHEA ASSET", size=14, color=TEAL, bold=True)
    add_text(slide, Inches(0.8), Inches(0.85), Inches(11), Inches(0.7),
             "A unique dataset purpose-built for ophthalmic AI",
             size=28, color=WHITE, bold=True)

    for i, (num, lbl) in enumerate([
        ("368K", "Subjects"), ("734K", "Eyes"),
        ("1.2M", "Images"), ("64", "Clinical Sites"),
    ]):
        x = 0.8 + i * 3.0
        add_text(slide, Inches(x), Inches(2.0), Inches(2.5), Inches(0.8),
                 num, size=48, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, Inches(x), Inches(2.75), Inches(2.5), Inches(0.4),
                 lbl, size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    for i, (title, body) in enumerate([
        ("Paired OCT + CFP", "Both modalities for every subject"),
        ("AI Models Pre-Computed", "RETscreenAI (fluid/RPE detection),\nAutoMorph (vessel analysis)"),
        ("US + Australia Sites", "64 primary care optometry\npractices across two countries"),
        ("Structured Analysis", "ETDRS grids, RNFL sectors,\ndisc topography, axial length"),
    ]):
        x = 0.8 + i * 3.0
        card = add_shape(slide, Inches(x), Inches(3.7), Inches(2.7), Inches(1.5), CARD_BG_DARK)
        add_shape(slide, Inches(x), Inches(3.7), Inches(2.7), Inches(0.05), TEAL)
        add_text(slide, Inches(x + 0.15), Inches(3.8), Inches(2.4), Inches(0.3),
                 title, size=13, color=TEAL, bold=True)
        add_text(slide, Inches(x + 0.15), Inches(4.15), Inches(2.4), Inches(0.9),
                 body, size=11, color=LIGHT_GRAY)

    # ── SLIDE 4: WHAT WE CAN PRE-SCREEN TODAY ──
    slide = prs.slides.add_slide(blank)
    add_shape(slide, Inches(0), Inches(0), prs.slide_width, prs.slide_height, OFFWHITE)
    add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), TEAL)
    add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.5),
             "WHAT WE CAN PRE-SCREEN TODAY", size=14, color=TEAL, bold=True)
    add_text(slide, Inches(0.8), Inches(0.85), Inches(11), Inches(0.7),
             f"{n_unique} trials mapped, {n_recruiting} currently recruiting",
             size=28, color=NAVY, bold=True)

    # Left: direct pre-screening fields
    add_text(slide, Inches(0.8), Inches(1.7), Inches(5), Inches(0.4),
             "Direct Pre-Screening (no external data needed)", size=14,
             color=NAVY, bold=True)
    for j, f in enumerate(direct_fields[:7]):
        add_text(slide, Inches(0.8), Inches(2.2 + j * 0.42), Inches(5.5), Inches(0.35),
                 f"\u2713  {f['field_name']}  \u2014  {f['direct_use'][:55]}",
                 size=11, color=DARK_GRAY)

    # Right: what still requires site data
    add_text(slide, Inches(7.0), Inches(1.7), Inches(5.5), Inches(0.4),
             "Requires Chart Review / Site Data", size=14,
             color=NAVY, bold=True)
    for j, c in enumerate(not_eval_criteria[:6]):
        add_text(slide, Inches(7.0), Inches(2.2 + j * 0.42), Inches(5.8), Inches(0.35),
                 f"\u2717  {c['criterion_text'][:55]}",
                 size=11, color=MID_GRAY)

    # Bottom: condition bar chart
    add_text(slide, Inches(0.8), Inches(5.0), Inches(8), Inches(0.3),
             "Trials by condition (includes cross-listed):", size=12,
             color=NAVY, bold=True)
    bar_y = 5.4
    max_count = max(cond_counts.values()) if cond_counts else 1
    for i, (cond, count) in enumerate(sorted(cond_counts.items(),
                                              key=lambda x: -x[1])[:6]):
        x = 0.8 + i * 2.0
        bar_w = (count / max_count) * 1.5
        add_shape(slide, Inches(x), Inches(bar_y), Inches(bar_w), Inches(0.3), TEAL)
        add_text(slide, Inches(x), Inches(bar_y + 0.3), Inches(1.8), Inches(0.5),
                 f"{condition_label(cond)}\n({count})", size=9, color=MID_GRAY,
                 align=PP_ALIGN.LEFT)

    # ── SLIDE 5: REVENUE MODEL ──
    slide = prs.slides.add_slide(blank)
    add_shape(slide, Inches(0), Inches(0), prs.slide_width, prs.slide_height, NAVY)
    add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), TEAL)
    add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.5),
             "REVENUE MODEL", size=14, color=TEAL, bold=True)
    add_text(slide, Inches(0.8), Inches(0.85), Inches(11), Inches(0.7),
             "Three revenue streams from the same data asset",
             size=28, color=WHITE, bold=True)

    for i, (title, price, unit, desc, timeline, color) in enumerate([
        ("Feasibility Reports", "$25K\u2013$75K", "per report",
         "Data-backed site selection and\npopulation sizing for specific\ntrial protocols.",
         "NOW", GREEN),
        ("Patient Pre-Screening", "$300\u2013$1,000", "per qualified referral",
         "Real-time anatomical pre-screening\nat point of care. Optometrist\nsees alert, patient gets referred.",
         "6\u201312 MO", TEAL),
        ("Real-World Evidence", "Platform", "subscription",
         "Longitudinal OCT + CFP data for\nnatural history studies and\npost-market evidence generation.",
         "12\u201324 MO", LIGHT_TEAL),
    ]):
        x = 0.8 + i * 4.0
        card = add_shape(slide, Inches(x), Inches(1.9), Inches(3.6), Inches(3.8), CARD_BG_DARK)
        add_shape(slide, Inches(x), Inches(1.9), Inches(3.6), Inches(0.06), color)
        badge = add_shape(slide, Inches(x + 2.4), Inches(2.1), Inches(1.0), Inches(0.35), color)
        tf = badge.text_frame
        tf.paragraphs[0].text = timeline
        tf.paragraphs[0].font.size = Pt(10)
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        add_text(slide, Inches(x + 0.2), Inches(2.1), Inches(2.2), Inches(0.4),
                 title, size=16, color=WHITE, bold=True)
        add_text(slide, Inches(x + 0.2), Inches(2.6), Inches(3.2), Inches(0.6),
                 price, size=36, color=color, bold=True)
        add_text(slide, Inches(x + 0.2), Inches(3.2), Inches(3.2), Inches(0.3),
                 unit, size=13, color=MID_GRAY)
        add_text(slide, Inches(x + 0.2), Inches(3.7), Inches(3.2), Inches(1.5),
                 desc, size=12, color=LIGHT_GRAY)

    # ── SLIDE 6: GTM ROADMAP ──
    slide = prs.slides.add_slide(blank)
    add_shape(slide, Inches(0), Inches(0), prs.slide_width, prs.slide_height, OFFWHITE)
    add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), TEAL)
    add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.5),
             "NEXT STEPS", size=14, color=TEAL, bold=True)
    add_text(slide, Inches(0.8), Inches(0.85), Inches(11), Inches(0.7),
             "From analysis to revenue in 12 months",
             size=28, color=NAVY, bold=True)

    steps = [
        ("1", "Quantify matchable population for top 10 recruiting trials",
         "Week 1\u20132"),
        ("2", "Secure commercial data rights and governance approval",
         "Week 2\u20134"),
        ("3", "Build rules engine MVP (CST + fluid + demographics)",
         "Week 3\u20136"),
        ("4", "Pilot feasibility reports with 2\u20133 CROs",
         "Week 4\u20138"),
        ("5", "Deliver paid engagements and build case studies",
         "Week 8\u201312"),
    ]
    for i, (num, desc, timeline) in enumerate(steps):
        y = 1.8 + i * 0.95
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
            Inches(0.8), Inches(y + 0.05), Inches(0.5), Inches(0.5))
        circle.fill.solid()
        circle.fill.fore_color.rgb = TEAL
        circle.line.fill.background()
        tf = circle.text_frame
        tf.paragraphs[0].text = num
        tf.paragraphs[0].font.size = Pt(20)
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        add_text(slide, Inches(1.5), Inches(y + 0.05), Inches(8), Inches(0.45),
                 desc, size=16, color=NAVY, bold=False)
        badge = add_shape(slide, Inches(10), Inches(y + 0.1),
                          Inches(1.5), Inches(0.35), TEAL)
        tf = badge.text_frame
        tf.paragraphs[0].text = timeline
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # ── SLIDE 7: CLOSING ──
    slide = prs.slides.add_slide(blank)
    add_shape(slide, Inches(0), Inches(0), prs.slide_width, prs.slide_height, NAVY)
    add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), TEAL)
    add_text(slide, Inches(1.2), Inches(2.0), Inches(11), Inches(1),
             "iDHEA can pre-screen for anatomical trial criteria\n"
             "across 64 sites, today.",
             size=32, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1.2), Inches(3.4), Inches(11), Inches(0.8),
             f"368K patients. {n_unique} mappable trials. {n_recruiting} recruiting now.\n"
             f"No new imaging capture needed for anatomy-only pre-screening.",
             size=20, color=TEAL, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1.2), Inches(4.6), Inches(11), Inches(0.6),
             "Full eligibility requires BCVA, IOP, HbA1c, and treatment history\n"
             "from chart review or site-level data integration.",
             size=14, color=MID_GRAY, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1.2), Inches(5.8), Inches(11), Inches(0.5),
             f"Confidential  |  {date.today().strftime('%B %Y')}  |  idhea.net",
             size=13, color=MID_GRAY, align=PP_ALIGN.CENTER)

    out_path = OUTPUTS / "gtm_pharma_deck.pptx"
    prs.save(str(out_path))
    print(f"Wrote {out_path}")


if __name__ == "__main__":
    generate()
